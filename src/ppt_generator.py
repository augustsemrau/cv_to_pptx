from pptx import Presentation
from pptx.util import Inches


def json_to_ppt(cv: dict) -> bytes:
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = cv.get("name", "")
    subtitle = title_slide.placeholders[1]
    subtitle.text = cv.get("title", "")

    # Profile slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Profile"
    slide.placeholders[1].text = cv.get("profile", "")

    # Education slide
    if cv.get("educations"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Education"
        body = slide.placeholders[1].text_frame
        for edu in cv["educations"]:
            p = body.add_paragraph()
            p.text = f"{edu['start_year']}-{edu['end_year']}: {edu['degree']} at {edu['institution']}"

    # Projects slide
    if cv.get("projects"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Projects"
        body = slide.placeholders[1].text_frame
        for proj in cv["projects"]:
            p = body.add_paragraph()
            p.text = f"{proj['name']}: {proj['description']}"

    # Serialize to bytes
    from io import BytesIO
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()
